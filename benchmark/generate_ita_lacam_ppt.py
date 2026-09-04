#!/usr/bin/env python3
"""Generate a visual PPT that explains ITA-LaCAM in the carrier project.

Each slide is first rendered as a 1920x1080 PNG and then embedded full-bleed
into the PPTX. This keeps Chinese typography and layout deterministic.
"""

from pathlib import Path
from urllib.request import urlretrieve
import math

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "benchmark" / "viz_web"
PPT_PATH = OUT_DIR / "ita_lacam_project_explainer.pptx"
PREVIEW_PATH = OUT_DIR / "ita_lacam_project_explainer_preview.png"
WORK_DIR = Path("/tmp/ita-lacam-ppt-build")
FONT_DIR = Path("/tmp/ita-lacam-ppt-fonts")

W, H = 1920, 1080
BG = "#07111f"
PANEL = "#0d1c2f"
PANEL_2 = "#10243a"
LINE = "#2b4662"
TEXT = "#f2f8ff"
MUTED = "#a9bed3"
BLUE = "#68c6ff"
GREEN = "#57e0a2"
YELLOW = "#ffd26f"
PINK = "#f58fc6"
VIOLET = "#b69cff"
RED = "#ff7f86"
WHITE = "#ffffff"

FONT_URLS = {
    "regular": "https://fonts.gstatic.com/s/notosanssc/v40/k3kCo84MPvpLmixcA63oeAL7Iqp5IZJF9bmaG9_FnYw.ttf",
    "bold": "https://fonts.gstatic.com/s/notosanssc/v40/k3kCo84MPvpLmixcA63oeAL7Iqp5IZJF9bmaGzjCnYw.ttf",
    "black": "https://fonts.gstatic.com/s/notosanssc/v40/k3kCo84MPvpLmixcA63oeAL7Iqp5IZJF9bmaG3bCnYw.ttf",
}


def rgb(value):
    value = value.lstrip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def ensure_fonts():
    FONT_DIR.mkdir(parents=True, exist_ok=True)
    paths = {}
    for name, url in FONT_URLS.items():
        path = FONT_DIR / ("NotoSansSC-%s.ttf" % name)
        if not path.exists():
            urlretrieve(url, str(path))
        paths[name] = path
    return paths


FONT_PATHS = ensure_fonts()


def font(size, weight="regular"):
    return ImageFont.truetype(str(FONT_PATHS[weight]), size)


def rounded(draw, box, radius=24, fill=PANEL, outline=None, width=2):
    draw.rounded_rectangle(box, radius=radius, fill=rgb(fill),
                           outline=rgb(outline) if outline else None, width=width)


def text_size(draw, value, fnt):
    box = draw.textbbox((0, 0), value, font=fnt)
    return box[2] - box[0], box[3] - box[1]


def draw_text(draw, xy, value, size, color=TEXT, weight="regular",
              anchor="la", stroke=0, stroke_fill=None):
    fnt = font(size, weight)
    draw.text(xy, value, font=fnt, fill=rgb(color), anchor=anchor,
              stroke_width=stroke,
              stroke_fill=rgb(stroke_fill) if stroke_fill else None)
    return fnt


def wrap_lines(draw, value, fnt, max_width):
    lines = []
    for paragraph in value.split("\n"):
        if paragraph == "":
            lines.append("")
            continue
        current = ""
        for ch in paragraph:
            candidate = current + ch
            if current and text_size(draw, candidate, fnt)[0] > max_width:
                lines.append(current.rstrip())
                current = ch.lstrip() if ch == " " else ch
            else:
                current = candidate
        if current:
            lines.append(current.rstrip())
    return lines


def paragraph(draw, box, value, size=30, color=MUTED, weight="regular",
              line_gap=12, max_lines=None, align="left"):
    x1, y1, x2, y2 = box
    fnt = font(size, weight)
    lines = wrap_lines(draw, value, fnt, x2 - x1)
    if max_lines:
        lines = lines[:max_lines]
    line_h = text_size(draw, "国Ag", fnt)[1] + line_gap
    y = y1
    for line in lines:
        if y + line_h > y2:
            break
        if align == "center":
            x = (x1 + x2) / 2
            anchor = "ma"
        elif align == "right":
            x = x2
            anchor = "ra"
        else:
            x = x1
            anchor = "la"
        draw.text((x, y), line, font=fnt, fill=rgb(color), anchor=anchor)
        y += line_h
    return y


def pill(draw, box, label, fill, color=BG, size=24, outline=None):
    rounded(draw, box, radius=min(24, max(8, (box[3] - box[1]) // 3)), fill=fill,
            outline=outline, width=2)
    draw_text(draw, ((box[0] + box[2]) / 2, (box[1] + box[3]) / 2),
              label, size, color, "bold", anchor="mm")


def arrow(draw, start, end, color=BLUE, width=8, head=18):
    draw.line([start, end], fill=rgb(color), width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    left = (end[0] - head * math.cos(angle - math.pi / 6),
            end[1] - head * math.sin(angle - math.pi / 6))
    right = (end[0] - head * math.cos(angle + math.pi / 6),
             end[1] - head * math.sin(angle + math.pi / 6))
    draw.polygon([end, left, right], fill=rgb(color))


def dashed_line(draw, start, end, color=LINE, width=5, dash=15, gap=10):
    dx, dy = end[0] - start[0], end[1] - start[1]
    length = math.hypot(dx, dy)
    if not length:
        return
    ux, uy = dx / length, dy / length
    pos = 0
    while pos < length:
        stop = min(length, pos + dash)
        draw.line([(start[0] + ux * pos, start[1] + uy * pos),
                   (start[0] + ux * stop, start[1] + uy * stop)],
                  fill=rgb(color), width=width)
        pos += dash + gap


def slide_base(index, title, kicker):
    img = Image.new("RGB", (W, H), rgb(BG))
    draw = ImageDraw.Draw(img)
    draw.ellipse((-330, -420, 760, 640), fill=rgb("#0b2741"))
    draw.ellipse((1450, -260, 2200, 500), fill=rgb("#191c43"))
    draw.rectangle((0, 0, W, 12), fill=rgb(BLUE))
    draw_text(draw, (96, 66), kicker.upper(), 22, BLUE, "bold")
    draw_text(draw, (96, 110), title, 56, TEXT, "black")
    draw.line((96, 190, 1824, 190), fill=rgb(LINE), width=2)
    draw_text(draw, (96, 1037), "ITA-LaCAM × Carrier-LaCAM", 20, MUTED, "regular")
    draw_text(draw, (1824, 1037), "%02d / 10" % index, 20, MUTED, "bold", anchor="ra")
    for i in range(10):
        x = 730 + i * 48
        draw.ellipse((x, 1027, x + 14, 1041),
                     fill=rgb(BLUE if i + 1 == index else LINE))
    return img, draw


def sticky(draw, box, title, line2, accent=YELLOW):
    rounded(draw, box, radius=22, fill="#282416", outline=accent, width=3)
    pill(draw, (box[2] - 134, box[1] + 18, box[2] - 20, box[1] + 56),
         "DERIVED", accent, BG, 17)
    draw_text(draw, (box[0] + 28, box[1] + 28), "guide(X)", 22, MUTED, "bold")
    draw_text(draw, (box[0] + 28, box[1] + 78), title, 35, accent, "black")
    draw_text(draw, (box[0] + 28, box[1] + 132), line2, 24, TEXT, "regular")


def node(draw, center, label, subtitle, active=True, color=BLUE, width=190):
    x, y = center
    box = (x - width // 2, y - 62, x + width // 2, y + 62)
    rounded(draw, box, 20, PANEL_2, color if active else LINE, 4 if active else 2)
    draw_text(draw, (x, y - 12), label, 34, TEXT if active else MUTED, "black", anchor="mm")
    draw_text(draw, (x, y + 34), subtitle, 19, MUTED, "regular", anchor="mm")
    return box


def draw_grid(draw, box, robot_x=4, selected="right"):
    x1, y1, x2, y2 = box
    cols, rows = 9, 5
    cw, ch = (x2 - x1) / cols, (y2 - y1) / rows
    rounded(draw, box, 22, "#081725", LINE, 3)
    for i in range(1, cols):
        draw.line((x1 + i * cw, y1, x1 + i * cw, y2), fill=rgb("#193049"), width=2)
    for i in range(1, rows):
        draw.line((x1, y1 + i * ch, x2, y1 + i * ch), fill=rgb("#193049"), width=2)

    def cell_center(cx, cy):
        return x1 + (cx + .5) * cw, y1 + (cy + .5) * ch

    def actor(cx, cy, label, fill, outline=None, circle=False, text_color=BG):
        px, py = cell_center(cx, cy)
        r = min(cw, ch) * .31
        shape = (px - r, py - r, px + r, py + r)
        if circle:
            draw.ellipse(shape, fill=rgb(fill), outline=rgb(outline) if outline else None, width=5)
        else:
            rounded(draw, shape, 14, fill, outline, 4 if outline else 2)
        draw_text(draw, (px, py), label, 27, text_color, "black", anchor="mm")
        return px, py

    gl = actor(1, 1, "G左", "#102437", GREEN, False, GREEN)
    gr = actor(7, 1, "G右", "#102437", PINK, False, PINK)
    b = actor(4, 1, "B", GREEN, None, False, BG)
    vl = actor(3, 2, "V左", "#7252ba", "#dacfff", False, WHITE)
    vr = actor(6, 2, "V右", "#7252ba", "#dacfff", False, WHITE)
    actor(4, 3, "空", "#15243a", VIOLET, False, VIOLET)
    r = actor(robot_x, 4, "R", "#e9f8ff", BLUE, True, BG)

    target_goal = gl if selected == "left" else gr
    color = GREEN if selected == "left" else PINK
    mid = ((b[0] + target_goal[0]) / 2, y1 + ch * .35)
    dashed_line(draw, b, mid, color, 5, 12, 8)
    dashed_line(draw, mid, target_goal, color, 5, 12, 8)
    chosen_leaf = vl if selected == "left" else vr
    draw.ellipse((chosen_leaf[0] - 43, chosen_leaf[1] - 43,
                  chosen_leaf[0] + 43, chosen_leaf[1] + 43),
                 outline=rgb(YELLOW), width=5)
    return {"B": b, "GL": gl, "GR": gr, "VL": vl, "VR": vr, "R": r}


def cost_row(draw, box, label, dep, robot, total, winner=False, color=GREEN):
    rounded(draw, box, 18, "#0b1b2d", color if winner else LINE, 3 if winner else 2)
    draw_text(draw, (box[0] + 24, box[1] + 24), label, 27, TEXT, "bold")
    draw_text(draw, (box[2] - 24, box[1] + 24), str(total), 34,
              color if winner else MUTED, "black", anchor="ra")
    y = box[1] + 72
    draw_text(draw, (box[0] + 24, y), "dependency %d" % dep, 21, VIOLET, "bold")
    draw_text(draw, (box[0] + 218, y), "+ robot %d" % robot, 21, BLUE, "bold")
    draw_text(draw, (box[0] + 390, y), "= %d" % total, 21,
              color if winner else MUTED, "black")
    bar_x1, bar_x2 = box[0] + 24, box[2] - 24
    bar_y = box[3] - 34
    draw.rounded_rectangle((bar_x1, bar_y, bar_x2, bar_y + 12), radius=6, fill=rgb("#06101c"))
    fill_x = bar_x1 + (bar_x2 - bar_x1) * total / 10
    draw.rounded_rectangle((bar_x1, bar_y, fill_x, bar_y + 12), radius=6,
                           fill=rgb(color if winner else "#617b95"))


def make_slides():
    slides = []

    # Slide 1
    img, d = slide_base(1, "ITA-LaCAM 到底在做什么？", "先说结论")
    paragraph(d, (96, 225, 1580, 330),
              "LaCAM 搜索物理状态 X；assignment 在每个 X 上重新计算，只负责告诉 PIBT“下一步先试什么”。",
              38, TEXT, "bold", 14)
    centers = [(330, 620), (960, 620), (1590, 620)]
    guides = [("τ(X0): B→G右", "ρ(X0): R→V右"),
              ("τ(X1): B→G左", "ρ(X1): R→V左"),
              ("τ(X2): 再重算", "物理变化，guide 也可变")]
    for i, center in enumerate(centers):
        if i:
            arrow(d, (centers[i - 1][0] + 120, 620), (center[0] - 120, 620), GREEN, 8, 22)
        node(d, center, "X%d" % i, "physical state", True, BLUE, 220)
        sticky(d, (center[0] - 190, 760, center[0] + 190, 940),
               guides[i][0], guides[i][1])
        dashed_line(d, (center[0], 684), (center[0], 760), YELLOW, 4, 10, 7)
    draw_text(d, (960, 390), "树里存 X；每个 X 旁边现算 guide(X)", 44, YELLOW, "black", anchor="ma")
    slides.append(img)

    # Slide 2
    img, d = slide_base(2, "先把两样东西彻底分开", "最重要的一页")
    rounded(d, (96, 245, 895, 920), 30, PANEL, BLUE, 4)
    pill(d, (132, 280, 310, 330), "真的状态", BLUE, BG, 22)
    draw_text(d, (132, 375), "物理状态 X", 48, TEXT, "black")
    paragraph(d, (132, 455, 835, 700),
              "机器人位置\n目标货架与匿名货架位置\nvacancy / carrying / custody",
              34, TEXT, "regular", 20)
    rounded(d, (132, 750, 835, 860), 18, "#0a1727", LINE, 2)
    draw_text(d, (166, 788), "进入 CLOSED key，也形成 LaCAM 树节点", 29, BLUE, "bold")

    rounded(d, (1025, 245, 1824, 920), 30, "#282416", YELLOW, 4)
    pill(d, (1061, 280, 1300, 330), "临时导航纸条", YELLOW, BG, 22)
    draw_text(d, (1061, 375), "guide(X)", 48, YELLOW, "black")
    paragraph(d, (1061, 455, 1764, 700),
              "τguide：target → goal\nDependency：当前 ready leaf\nρ：robot → task",
              34, TEXT, "regular", 20)
    rounded(d, (1061, 750, 1764, 860), 18, "#17150e", "#7f6a31", 2)
    draw_text(d, (1095, 788), "每个新 X 都重算；不进入 CLOSED key", 29, YELLOW, "bold")
    draw_text(d, (960, 960), "CLOSED key = X        guide(X) = derived", 30, MUTED, "bold", anchor="ma")
    slides.append(img)

    # Slide 3
    img, d = slide_base(3, "用一个最小例子：B 有两个合法 Goal", "问题设置")
    draw_grid(d, (96, 245, 1260, 895), robot_x=5, selected="right")
    rounded(d, (1325, 245, 1824, 895), 28, PANEL, LINE, 3)
    draw_text(d, (1370, 295), "当前物理节点：X0", 34, BLUE, "black")
    paragraph(d, (1370, 375, 1775, 560),
              "G左、G右都合法。\nV左、V右分别是两个方向当前真正可执行的 dependency leaf。",
              28, TEXT, "regular", 14)
    d.line((1370, 610, 1775, 610), fill=rgb(LINE), width=2)
    draw_text(d, (1370, 660), "现在要回答两件事", 27, MUTED, "bold")
    draw_text(d, (1370, 720), "① B 暂时朝哪个 Goal？", 31, YELLOW, "bold")
    draw_text(d, (1370, 780), "② R 暂时执行哪个 leaf？", 31, VIOLET, "bold")
    draw_text(d, (960, 945), "关键：必须把两个候选都评价完，不能先固定一个 Goal 再说。", 34, GREEN, "black", anchor="ma")
    slides.append(img)

    # Slide 4
    img, d = slide_base(4, "在 X0 上：先评价所有 (B, Goal)", "每节点动态 assignment")
    rounded(d, (96, 245, 960, 895), 28, PANEL, LINE, 3)
    draw_text(d, (138, 290), "候选 1：B → G左", 31, GREEN, "bold")
    paragraph(d, (138, 350, 900, 470),
              "Dependency：B→G左 ← blocker ← V左（ready）",
              26, TEXT, "regular", 10)
    cost_row(d, (138, 500, 900, 650), "G左 当前执行代价", 4, 3, 7, False, GREEN)
    draw_text(d, (138, 720), "候选 2：B → G右", 31, PINK, "bold")
    paragraph(d, (138, 770, 900, 850),
              "Dependency：B→G右 ← V右（ready）",
              26, TEXT, "regular", 10)

    rounded(d, (1020, 245, 1824, 895), 28, "#171827", PINK, 4)
    cost_row(d, (1062, 300, 1782, 465), "G右 当前执行代价", 4, 2, 6, True, PINK)
    draw_text(d, (1422, 545), "6 < 7", 66, YELLOW, "black", anchor="mm")
    sticky(d, (1110, 625, 1734, 825), "τguide(X0): B → G右",
           "ρ(X0): R → V右（ready leaf）")
    draw_text(d, (960, 950), "这不是“固定右目标”；只是 X0 的 preferred guide。", 34, TEXT, "black", anchor="ma")
    slides.append(img)

    # Slide 5
    img, d = slide_base(5, "PIBT 和 LaCAM 分工不同", "谁负责什么")
    rounded(d, (96, 275, 790, 865), 30, "#282416", YELLOW, 4)
    draw_text(d, (145, 335), "Carrier-PIBT", 45, YELLOW, "black")
    draw_text(d, (145, 415), "按 guide 生成第一个 successor", 31, TEXT, "bold")
    paragraph(d, (145, 500, 735, 700),
              "例如：\nτ 选 G右，ρ 选 V右\n→ PIBT 首先提出 MoveRight",
              31, MUTED, "regular", 16)
    pill(d, (145, 750, 560, 815), "快，但可能选错", RED, WHITE, 27)

    arrow(d, (820, 570), (1050, 570), BLUE, 10, 26)

    rounded(d, (1080, 275, 1824, 865), 30, PANEL, GREEN, 4)
    draw_text(d, (1130, 335), "LaCAM local constraint tree", 42, GREEN, "black")
    draw_text(d, (1130, 415), "保存并惰性枚举其他 primitive choices", 30, TEXT, "bold")
    paragraph(d, (1130, 500, 1765, 710),
              "首选走坏时：\n回到同一个 X0\n加一条 local constraint\n再次调用 PIBT 补全另一个 successor",
              29, MUTED, "regular", 14)
    pill(d, (1130, 750, 1680, 815), "preferred ≠ only", GREEN, BG, 27)
    draw_text(d, (960, 935), "PIBT 决定“先试谁”；LaCAM 决定“选错后还能继续试”。", 37, BLUE, "black", anchor="ma")
    slides.append(img)

    # Slide 6
    img, d = slide_base(6, "如果右边 Guide 走坏，会发生什么？", "LaCAM 兜底")
    node(d, (620, 355), "X0", "physical state", True, BLUE, 230)
    dashed_line(d, (565, 420), (365, 600), RED, 7, 16, 10)
    dashed_line(d, (675, 420), (900, 600), GREEN, 7, 16, 10)
    node(d, (330, 670), "Xloop", "duplicate", True, RED, 250)
    node(d, (930, 670), "X1", "alternative successor", True, GREEN, 290)
    pill(d, (550, 510, 930, 575), "local constraint: R ≠ right", VIOLET, WHITE, 23)
    rounded(d, (1130, 270, 1824, 875), 28, PANEL, LINE, 3)
    draw_text(d, (1180, 320), "注意：没有发生这些事", 34, RED, "black")
    paragraph(d, (1180, 405, 1760, 640),
              "✕ 没有创建 Task-LaCAM\n✕ 没有把 G右 固定后重跑整条路径\n✕ 没有把 assignment 放进 CLOSED",
              30, TEXT, "regular", 18)
    d.line((1180, 700, 1760, 700), fill=rgb(LINE), width=2)
    paragraph(d, (1180, 740, 1760, 840),
              "只是在 X0 内禁止这一次右动作，然后生成另一个物理 successor。",
              29, GREEN, "bold", 12)
    draw_text(d, (650, 930), "错误 guide 最多浪费搜索顺序，不会删除其他合法动作。", 35, YELLOW, "black", anchor="ma")
    slides.append(img)

    # Slide 7
    img, d = slide_base(7, "到了 X1：assignment 立即重新计算", "ITA 的核心")
    rounded(d, (96, 250, 900, 890), 28, PANEL, LINE, 3)
    draw_text(d, (140, 295), "在 X0", 36, BLUE, "black")
    draw_text(d, (140, 355), "R 更靠近 V右", 29, MUTED, "bold")
    cost_row(d, (140, 430, 850, 580), "B → G左", 4, 3, 7, False, GREEN)
    cost_row(d, (140, 620, 850, 770), "B → G右", 4, 2, 6, True, PINK)
    draw_text(d, (495, 825), "τ(X0): B→G右", 34, YELLOW, "black", anchor="ma")

    arrow(d, (920, 575), (1000, 575), BLUE, 10, 24)
    pill(d, (895, 475, 1025, 535), "R 左移", BLUE, BG, 22)

    rounded(d, (1020, 250, 1824, 890), 28, "#10251f", GREEN, 4)
    draw_text(d, (1065, 295), "在 X1", 36, GREEN, "black")
    draw_text(d, (1065, 355), "物理位置变了：R 更靠近 V左", 29, TEXT, "bold")
    cost_row(d, (1065, 430, 1775, 580), "B → G左", 4, 1, 5, True, GREEN)
    cost_row(d, (1065, 620, 1775, 770), "B → G右", 4, 3, 7, False, PINK)
    draw_text(d, (1420, 825), "τ(X1): B→G左", 34, YELLOW, "black", anchor="ma")
    draw_text(d, (960, 950), "不是外层“重跑规划”；X1 创建时就直接得到新的 guide(X1)。", 35, TEXT, "black", anchor="ma")
    slides.append(img)

    # Slide 8
    img, d = slide_base(8, "映射回我们的项目：始终只有一个循环", "正确架构")
    labels = [
        ("物理节点 X", "robot / shelf / vacancy / carrying", BLUE),
        ("编译所有 (b,g)", "每个候选 goal 的 dependency frontier", VIOLET),
        ("execution-aware cost", "blocker + vacancy + robot realization", PINK),
        ("τguide(X)", "Hungarian：target → goal", YELLOW),
        ("ready task pool", "只发射当前 dependency leaves", VIOLET),
        ("ρ(X)", "free robot → task", YELLOW),
        ("Carrier-PIBT", "先生成 preferred successor", BLUE),
        ("LaCAM constraints", "走坏则枚举其他 primitive actions", GREEN),
    ]
    xs = [96, 530, 964, 1398]
    ys = [285, 610]
    boxes = []
    for i, (title, sub, color) in enumerate(labels):
        row, col = divmod(i, 4)
        x, y = xs[col], ys[row]
        box = (x, y, x + 350, y + 210)
        boxes.append(box)
        rounded(d, box, 24, PANEL, color, 3)
        pill(d, (x + 24, y + 22, x + 104, y + 62), str(i + 1), color, BG, 20)
        draw_text(d, (x + 24, y + 87), title, 29, TEXT, "black")
        paragraph(d, (x + 24, y + 137, x + 326, y + 195), sub, 20, MUTED, "regular", 7)
        if col < 3:
            arrow(d, (x + 360, y + 105), (xs[col + 1] - 12, y + 105), BLUE, 6, 18)
    arrow(d, (1573, 515), (1573, 598), BLUE, 6, 18)
    arrow(d, (1390, 715), (1325, 715), GREEN, 6, 18)
    draw_text(d, (960, 930), "得到 X′ 后回到第 1 步：dependency、τ、ρ 全部从新物理状态重新计算。", 34, GREEN, "black", anchor="ma")
    slides.append(img)

    # Slide 9
    img, d = slide_base(9, "这和当前 v4.1 的关键差别", "项目决策")
    rounded(d, (96, 250, 900, 890), 28, "#25161c", RED, 4)
    pill(d, (135, 285, 360, 340), "当前 v4.1", RED, WHITE, 24)
    paragraph(d, (140, 395, 850, 760),
              "1. 先用 shelf-to-goal LB 得到 tau0\n\n2. 只为 tau0 已选的 Goal 编译 task\n\n3. 再做一次 execution-price repair\n\n4. 用 farthest / aging / 保留槽强推任务",
              29, TEXT, "regular", 16)
    draw_text(d, (500, 825), "问题：先选 Goal，后看真实执行结构", 29, RED, "black", anchor="ma")

    rounded(d, (1020, 250, 1824, 890), 28, "#10251f", GREEN, 4)
    pill(d, (1060, 285, 1345, 340), "更像 ITA-LaCAM", GREEN, BG, 24)
    paragraph(d, (1065, 395, 1770, 760),
              "1. 对每个候选 (target, goal) 先编译 dependency\n\n2. 把 vacancy 与 robot realization 算进 Cguide\n\n3. 每个 X 动态求 τguide(X)\n\n4. 由选中 Goal 的 ready leaves 动态求 ρ(X)",
              29, TEXT, "regular", 16)
    draw_text(d, (1420, 825), "结果：Goal 与 Task 是同一节点的联合 guide", 29, GREEN, "black", anchor="ma")
    draw_text(d, (960, 950), "τLB 仍单独用于 admissible h；只有 Lift/carry/drop custody 是硬承诺。", 33, YELLOW, "black", anchor="ma")
    slides.append(img)

    # Slide 10
    img, d = slide_base(10, "最后只记住这一句话", "总结")
    rounded(d, (170, 260, 1750, 540), 36, "#15283b", BLUE, 5)
    draw_text(d, (960, 340), "Goal / dependency / task 都是 guide(X)", 52, YELLOW, "black", anchor="ma")
    draw_text(d, (960, 430), "它们随物理节点动态变化，但不是新的搜索状态。", 45, TEXT, "black", anchor="ma")
    cards = [
        ("不是外层循环", "不是“固定 assignment → 求完整路径 → 再重分配”"),
        ("不是第二棵树", "不需要额外 Goal-LaCAM 或 Task-LaCAM"),
        ("不是强制调度", "不需要让最远 mission 永久抢占唯一槽"),
    ]
    for i, (title, sub) in enumerate(cards):
        x = 170 + i * 530
        rounded(d, (x, 620, x + 480, 875), 26, PANEL, [PINK, VIOLET, GREEN][i], 3)
        draw_text(d, (x + 30, 665), title, 32, [PINK, VIOLET, GREEN][i], "black")
        paragraph(d, (x + 30, 730, x + 450, 840), sub, 26, TEXT, "regular", 12)
    draw_text(d, (960, 945), "一棵 Carrier-LaCAM 物理树 + 每节点动态 assignment guidance", 37, BLUE, "black", anchor="ma")
    slides.append(img)

    return slides


def save_deck(slides):
    WORK_DIR.mkdir(parents=True, exist_ok=True)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    slide_paths = []
    for i, image in enumerate(slides, 1):
        path = WORK_DIR / ("slide-%02d.png" % i)
        image.save(str(path), "PNG", optimize=True)
        slide_paths.append(path)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    while len(prs.slides):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    for path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
    prs.core_properties.title = "ITA-LaCAM 与 Carrier-LaCAM：动态 Goal、Dependency、Task"
    prs.core_properties.subject = "One physical LaCAM tree with per-node assignment guidance"
    prs.core_properties.author = "Codex"
    prs.save(str(PPT_PATH))

    thumb_w, thumb_h = 640, 360
    sheet = Image.new("RGB", (thumb_w * 2 + 48, thumb_h * 5 + 72), rgb("#030913"))
    for i, image in enumerate(slides):
        thumb = image.resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        x = 16 + (i % 2) * (thumb_w + 16)
        y = 16 + (i // 2) * (thumb_h + 10)
        sheet.paste(thumb, (x, y))
    sheet.save(str(PREVIEW_PATH), "PNG", optimize=True)
    return slide_paths


if __name__ == "__main__":
    generated = save_deck(make_slides())
    print("PPT:", PPT_PATH)
    print("Preview:", PREVIEW_PATH)
    print("Slides:", len(generated))
